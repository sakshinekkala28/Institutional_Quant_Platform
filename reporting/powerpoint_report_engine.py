from __future__ import annotations

import logging

from pathlib import Path

import pandas as pd

from pptx import Presentation

from pptx.util import Inches

from core.settings import settings

# ==========================================================
# LOGGING
# ==========================================================

logging.basicConfig(

    level=logging.INFO,

    format=(

        "%(asctime)s | "

        "%(levelname)s | "

        "%(message)s"

    )

)

logger = logging.getLogger(

    __name__

)

# ==========================================================
# PATHS
# ==========================================================

ROOT_DIR = (

    settings
    .environment
    .ROOT_DIR

)

PERFORMANCE_DIR = (

    ROOT_DIR

    / "data"

    / "performance"

)

ALERT_GOVERNANCE_FILE = (

    ROOT_DIR

    / "data"

    / "monitoring"

    / "alert_governance.csv"

)

REPORT_DIR = (

    ROOT_DIR

    / "data"

    / "reports"

)

REPORT_DIR.mkdir(

    parents=True,

    exist_ok=True

)

FORECAST_FILE = (

    PERFORMANCE_DIR

    / "performance_forecast.csv"

)

GOVERNANCE_FILE = (

    PERFORMANCE_DIR

    / "governance_decision.csv"

)

SCENARIO_FILE = (

    PERFORMANCE_DIR

    / "scenario_analysis.csv"

)

STRESS_FILE = (

    PERFORMANCE_DIR

    / "stress_test_results.csv"

)

COMMITTEE_PACK_FILE = (

    PERFORMANCE_DIR

    / "investment_committee_pack.csv"

)

EXECUTIVE_FILE = (

    PERFORMANCE_DIR

    / "executive_dashboard.csv"

)

PPT_FILE = (

    REPORT_DIR

    / "investment_committee_presentation.pptx"

)

# ==========================================================
# REPOSITORY
# ==========================================================

class PPTRepository:

    @staticmethod
    def load_forecast():

        return pd.read_csv(

            FORECAST_FILE

        )

    @staticmethod
    def load_governance():

        return pd.read_csv(

            GOVERNANCE_FILE

        )

    @staticmethod
    def load_scenario():

        return pd.read_csv(

            SCENARIO_FILE

        )

    @staticmethod
    def load_stress():

        return pd.read_csv(

            STRESS_FILE

        )

    @staticmethod
    def load_committee():

        return pd.read_csv(

            COMMITTEE_PACK_FILE

        )
    
    @staticmethod
    def load_alert_governance():

        return pd.read_csv(

            ALERT_GOVERNANCE_FILE

        )

    @staticmethod
    def load_executive():

        return pd.read_csv(

            EXECUTIVE_FILE

        )
    
# ==========================================================
# VALIDATOR
# ==========================================================

class PPTValidator:

    @staticmethod
    def validate(

        forecast,

        governance

    ):

        if forecast.empty:

            raise ValueError(

                "Forecast Empty"

            )

        if governance.empty:

            raise ValueError(

                "Governance Empty"

            )
        
# ==========================================================
# COVER SLIDE
# ==========================================================

class CoverSlideBuilder:

    @staticmethod
    def build(

        prs,

        forecast,

        governance

    ):

        logger.info(

            "Building Cover Slide"

        )

        slide = (

            prs.slides.add_slide(

                prs.slide_layouts[0]

            )

        )

        forecast_row = (

            forecast

            .iloc[0]

        )

        governance_row = (

            governance

            .iloc[0]

        )

        slide.shapes.title.text = (

            "Institutional Investment Committee Review"

        )

        slide.placeholders[1].text = (

            f"Forecast Regime: "

            f"{forecast_row['Forecast_Regime']}\n"

            f"Recommendation: "

            f"{forecast_row['Forecast_Recommendation']}\n"

            f"Decision: "

            f"{governance_row['Decision']}"

        )

# ==========================================================
# EXECUTIVE SLIDE
# ==========================================================

class ExecutiveSlideBuilder:

    @staticmethod
    def build(

        prs,

        executive

    ):

        logger.info(

            "Building Executive Slide"

        )

        slide = (

            prs.slides.add_slide(

                prs.slide_layouts[1]

            )

        )

        slide.shapes.title.text = (

            "Executive Dashboard"

        )

        textbox = (

            slide.shapes.add_textbox(

                Inches(0.5),

                Inches(1.2),

                Inches(8),

                Inches(4)

            )

        )

        text_frame = (

            textbox

            .text_frame

        )

        for _, row in executive.iterrows():

            text_frame.add_paragraph().text = (

                f"{row['Metric']}: "

                f"{row['Value']}"

            )

# ==========================================================
# FORECAST SLIDE
# ==========================================================

class ForecastSlideBuilder:

    @staticmethod
    def build(

        prs,

        forecast

    ):

        logger.info(

            "Building Forecast Slide"

        )

        slide = (

            prs.slides.add_slide(

                prs.slide_layouts[1]

            )

        )

        slide.shapes.title.text = (

            "Performance Forecast"

        )

        row = (

            forecast

            .iloc[0]

        )

        textbox = (

            slide.shapes.add_textbox(

                Inches(0.5),

                Inches(1.2),

                Inches(8),

                Inches(4)

            )

        )

        tf = (

            textbox

            .text_frame

        )

        tf.add_paragraph().text = (

            f"Expected Return (1M): "

            f"{row['Expected_Return_1M']:.2%}"

        )

        tf.add_paragraph().text = (

            f"Expected Return (3M): "

            f"{row['Expected_Return_3M']:.2%}"

        )

        tf.add_paragraph().text = (

            f"Expected Return (12M): "

            f"{row['Expected_Return_12M']:.2%}"

        )

        tf.add_paragraph().text = (

            f"Forecast Confidence: "

            f"{row['Forecast_Confidence']:.2f}"

        )

        tf.add_paragraph().text = (

            f"Forecast Risk Grade: "

            f"{row['Forecast_Risk_Grade']}"

        )

        tf.add_paragraph().text = (

            f"Forecast Regime: "

            f"{row['Forecast_Regime']}"

        )

# ==========================================================
# GOVERNANCE SLIDE
# ==========================================================

class GovernanceSlideBuilder:

    @staticmethod
    def build(

        prs,

        governance

    ):

        logger.info(

            "Building Governance Slide"

        )

        slide = (

            prs.slides.add_slide(

                prs.slide_layouts[1]

            )

        )

        slide.shapes.title.text = (

            "Governance Decision"

        )

        row = (

            governance

            .iloc[0]

        )

        textbox = (

            slide.shapes.add_textbox(

                Inches(0.5),

                Inches(1.2),

                Inches(8),

                Inches(4)

            )

        )

        tf = (

            textbox

            .text_frame

        )

        tf.add_paragraph().text = (

            f"Committee Score: "

            f"{row['Committee_Score']:.2f}"

        )

        tf.add_paragraph().text = (

            f"Governance Score: "

            f"{row['Governance_Score']:.2f}"

        )

        tf.add_paragraph().text = (

            f"Risk Score: "

            f"{row['Risk_Score']:.2f}"

        )

        tf.add_paragraph().text = (

            f"Decision: "

            f"{row['Decision']}"

        )

# ==========================================================
# SCENARIO ANALYSIS SLIDE
# ==========================================================

class ScenarioSlideBuilder:

    @staticmethod
    def build(

        prs,

        scenario

    ):

        logger.info(

            "Building Scenario Slide"

        )

        slide = (

            prs.slides.add_slide(

                prs.slide_layouts[1]

            )

        )

        slide.shapes.title.text = (

            "Scenario Analysis"

        )

        textbox = (

            slide.shapes.add_textbox(

                Inches(0.5),

                Inches(1.2),

                Inches(8.5),

                Inches(4)

            )

        )

        tf = (

            textbox

            .text_frame

        )

        ranked = (

            scenario

            .sort_values(

                "Scenario_Rank"

            )

        )

        for _, row in ranked.iterrows():

            tf.add_paragraph().text = (

                f"{row['Scenario']} | "

                f"Score: {row['Scenario_Score']:.2f} | "

                f"Rank: {row['Scenario_Rank']}"

            )

# ==========================================================
# STRESS TESTING SLIDE
# ==========================================================

class StressSlideBuilder:

    @staticmethod
    def build(

        prs,

        stress

    ):

        logger.info(

            "Building Stress Slide"

        )

        slide = (

            prs.slides.add_slide(

                prs.slide_layouts[1]

            )

        )

        slide.shapes.title.text = (

            "Stress Testing"

        )

        textbox = (

            slide.shapes.add_textbox(

                Inches(0.5),

                Inches(1.2),

                Inches(8.5),

                Inches(4)

            )

        )

        tf = (

            textbox

            .text_frame

        )

        ranked = (

            stress

            .sort_values(

                "Stress_Rank"

            )

        )

        for _, row in ranked.iterrows():

            tf.add_paragraph().text = (

                f"{row['Scenario']} | "

                f"Score: {row['Stress_Score']:.2f} | "

                f"Impact: {row['Return_Impact']:.2%}"

            )

# ==========================================================
# COMMITTEE PACK SLIDE
# ==========================================================

class CommitteePackSlideBuilder:

    @staticmethod
    def build(

        prs,

        committee

    ):

        logger.info(

            "Building Committee Pack Slide"

        )

        slide = (

            prs.slides.add_slide(

                prs.slide_layouts[1]

            )

        )

        slide.shapes.title.text = (

            "Investment Committee Pack"

        )

        textbox = (

            slide.shapes.add_textbox(

                Inches(0.3),

                Inches(1.0),

                Inches(9),

                Inches(5)

            )

        )

        tf = (

            textbox

            .text_frame

        )

        current_section = None

        for _, row in committee.iterrows():

            section = (

                row.get(

                    "Section",

                    "General"

                )

            )

            if section != current_section:

                current_section = section

                tf.add_paragraph().text = (

                    f"=== {section} ==="

                )

            tf.add_paragraph().text = (

                f"{row['Metric']}: "

                f"{row['Value']}"

            )



class AlertGovernanceSlideBuilder:

    @staticmethod
    def build(

        prs,

        alert_governance

    ):

        logger.info(

            "Building Alert Governance Slide"

        )

        slide = (

            prs.slides.add_slide(

                prs.slide_layouts[1]

            )

        )

        slide.shapes.title.text = (

            "Alert Governance"

        )

        metrics = dict(

            zip(

                alert_governance["Metric"],

                alert_governance["Value"]

            )

        )

        textbox = (

            slide.shapes.add_textbox(

                Inches(0.5),

                Inches(1.2),

                Inches(8),

                Inches(4)

            )

        )

        tf = (

            textbox.text_frame

        )

        tf.add_paragraph().text = (

            f"Health Score: "

            f"{metrics.get('Alert_Health_Score')}"

        )

        tf.add_paragraph().text = (

            f"Escalation: "

            f"{metrics.get('Alert_Escalation')}"

        )

        tf.add_paragraph().text = (

            f"Governance View: "

            f"{metrics.get('Governance_View')}"

        )

        tf.add_paragraph().text = (

            f"Trend Direction: "

            f"{metrics.get('Trend_Direction')}"

        )

        tf.add_paragraph().text = (

            f"Top Category: "

            f"{metrics.get('Top_Category')}"

        )


# ==========================================================
# SURVEILLANCE SLIDE
# ==========================================================

class SurveillanceSlideBuilder:

    @staticmethod
    def build(

        prs,

        executive

    ):

        logger.info(

            "Building Surveillance Slide"

        )

        slide = (

            prs.slides.add_slide(

                prs.slide_layouts[1]

            )

        )

        slide.shapes.title.text = (

            "Portfolio Monitoring"

        )

        textbox = (

            slide.shapes.add_textbox(

                Inches(0.5),

                Inches(1.2),

                Inches(8),

                Inches(4)

            )

        )

        tf = (

            textbox

            .text_frame

        )

        for _, row in executive.iterrows():

            metric = str(

                row["Metric"]

            )

            if (

                "Alert" in metric

                or

                "Health" in metric

            ):

                tf.add_paragraph().text = (

                    f"{metric}: "

                    f"{row['Value']}"

                )

# ==========================================================
# PPT BUILDER
# ==========================================================

class PPTBuilder:

    def build(

        self

    ):

        logger.info(

            "Building Presentation"

        )

        forecast = (

            PPTRepository

            .load_forecast()

        )

        governance = (

            PPTRepository

            .load_governance()

        )

        scenario = (

            PPTRepository

            .load_scenario()

        )

        stress = (

            PPTRepository

            .load_stress()

        )

        committee = (

            PPTRepository

            .load_committee()

        )

        alert_governance = (

            PPTRepository

            .load_alert_governance()

        )

        executive = (

            PPTRepository

            .load_executive()

        )

        PPTValidator.validate(

            forecast,

            governance

        )

        prs = Presentation()

        CoverSlideBuilder.build(

            prs,

            forecast,

            governance

        )

        ExecutiveSlideBuilder.build(

            prs,

            executive

        )

        ForecastSlideBuilder.build(

            prs,

            forecast

        )

        GovernanceSlideBuilder.build(

            prs,

            governance

        )

        ScenarioSlideBuilder.build(

            prs,

            scenario

        )

        StressSlideBuilder.build(

            prs,

            stress

        )

        CommitteePackSlideBuilder.build(

            prs,

            committee

        )

        AlertGovernanceSlideBuilder.build(

            prs,

            alert_governance

        )

        SurveillanceSlideBuilder.build(

            prs,

            executive

        )

        return prs
    

# ==========================================================
# PPT EXPORT ENGINE
# ==========================================================

class PPTExportEngine:

    @staticmethod
    def export():

        logger.info(

            "Exporting PowerPoint"

        )

        presentation = (

            PPTBuilder()

            .build()

        )

        presentation.save(

            str(

                PPT_FILE

            )

        )

        logger.info(

            "PowerPoint Exported"

        )

        return PPT_FILE
    
# ==========================================================
# EXECUTIVE SUMMARY
# ==========================================================

class PPTSummaryEngine:

    @staticmethod
    def build():

        executive = (

            PPTRepository

            .load_executive()

        )

        metrics = {}

        for _, row in executive.iterrows():

            metrics[

                row["Metric"]

            ] = row["Value"]

        return metrics
    
# ==========================================================
# RUNNER
# ==========================================================

def run_example():

    ppt_path = (

        PPTExportEngine

        .export()

    )

    summary = (

        PPTSummaryEngine

        .build()

    )

    print()

    print(

        "=" * 80

    )

    print(

        "INVESTMENT COMMITTEE PRESENTATION"

    )

    print(

        "=" * 80

    )

    print(

        f"Presentation: {ppt_path}"

    )

    print()

    for key, value in summary.items():

        print(

            f"{key}: {value}"

        )

    print(

        "=" * 80

    )

# ==========================================================
# MAIN
# ==========================================================

if __name__ == "__main__":

    run_example()